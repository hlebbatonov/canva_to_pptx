from pptx import Presentation
n = int(input()) #input how many slides you have
root = Presentation()
slide_layout = root.slide_layouts[6] #creating blenk slide
try:
    for i in range(1, n + 1): #iterate all files
        slide = root.slides.add_slide(slide_layout)
        left = top = 0
        pic = slide.shapes.add_picture(f'{i}.png', left - 0.1 * root.slide_width, top, height=root.slide_height)
except Exception as e:
    print(e)
root.save('done.pptx')